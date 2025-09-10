"""
Microsoft Office document metadata extractor for Word, Excel, PowerPoint
"""

import asyncio
from pathlib import Path
from typing import List, Optional, Dict, Any
import time

from .base_extractor import BaseExtractor, ExtractedMetadata


class OfficeExtractor(BaseExtractor):
    """Extract metadata from Microsoft Office documents"""
    
    def __init__(self):
        super().__init__()
        self.supported_types = [
            # Legacy Office formats
            'application/msword',
            'application/vnd.ms-excel',
            'application/vnd.ms-powerpoint',
            
            # Modern Office Open XML formats  
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            
            # OpenDocument formats
            'application/vnd.oasis.opendocument.text',
            'application/vnd.oasis.opendocument.spreadsheet',
            'application/vnd.oasis.opendocument.presentation'
        ]
    
    async def can_extract(self, file_path: Path, mime_type: str) -> bool:
        """Check if this extractor can handle Office documents"""
        return mime_type in self.supported_types or file_path.suffix.lower() in [
            '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.odt', '.ods', '.odp'
        ]
    
    async def extract_metadata(self, file_path: Path, mime_type: str) -> ExtractedMetadata:
        """Extract comprehensive metadata from Office documents"""
        start_time = time.time()
        metadata = ExtractedMetadata()
        metadata.extractor_version = f"OfficeExtractor-{self.version}"
        
        try:
            # Run Office extraction in thread pool to avoid blocking
            loop = asyncio.get_event_loop()
            office_data = await loop.run_in_executor(None, self._extract_office_sync, file_path, mime_type)
            
            if office_data:
                # Core document properties
                metadata.title = office_data.get('title')
                metadata.author = office_data.get('author')
                metadata.subject = office_data.get('subject')
                metadata.creator = office_data.get('creator')
                metadata.keywords = office_data.get('keywords', [])
                metadata.creation_date = office_data.get('creation_date')
                metadata.modification_date = office_data.get('modification_date')
                
                # Office-specific metadata
                metadata.page_count = office_data.get('page_count')
                metadata.word_count = office_data.get('word_count')
                metadata.character_count = office_data.get('character_count')
                metadata.language = office_data.get('language')
                
                # Content
                metadata.text_content = office_data.get('text_content', '')
                
                # Document structure
                metadata.tables_detected = office_data.get('tables_count', 0)
                metadata.images_detected = office_data.get('images_count', 0)
                metadata.links_detected = office_data.get('links_count', 0)
                
                # Security
                metadata.password_protected = office_data.get('password_protected', False)
                
                # Raw metadata
                metadata.raw_metadata = office_data.get('raw_metadata', {})
                
                # Process extracted text
                if metadata.text_content:
                    if not metadata.word_count:  # If not extracted from document properties
                        metadata.word_count = self._count_words(metadata.text_content)
                    if not metadata.character_count:
                        metadata.character_count = self._count_characters(metadata.text_content)
                    if not metadata.language:
                        metadata.language = self._detect_language(metadata.text_content)
                    
                    # Extract keywords if not already present
                    if not metadata.keywords:
                        metadata.keywords = self._extract_keywords(metadata.text_content)
                    
                    # Text extraction confidence
                    metadata.text_extraction_confidence = self._assess_extraction_quality(
                        metadata.text_content, metadata.word_count
                    )
                    
                    # Content-based suggestions
                    metadata.suggested_categories = self._suggest_categories(
                        metadata.text_content, mime_type, office_data
                    )
                    metadata.suggested_tags = self._suggest_tags(
                        metadata.text_content, metadata.title, mime_type
                    )
                
        except Exception as e:
            error_msg = f"Office extraction failed: {str(e)}"
            metadata.errors.append(error_msg)
            self.logger.error(error_msg)
        
        # Record processing time
        end_time = time.time()
        metadata.extraction_duration_ms = int((end_time - start_time) * 1000)
        
        return metadata
    
    def _extract_office_sync(self, file_path: Path, mime_type: str) -> Dict[str, Any]:
        """Synchronous Office extraction (runs in thread pool)"""
        suffix = file_path.suffix.lower()
        
        # Determine document type and extract accordingly
        if suffix in ['.docx', '.doc'] or 'wordprocessingml' in mime_type:
            return self._extract_word_document(file_path)
        elif suffix in ['.xlsx', '.xls'] or 'spreadsheetml' in mime_type:
            return self._extract_excel_document(file_path)
        elif suffix in ['.pptx', '.ppt'] or 'presentationml' in mime_type:
            return self._extract_powerpoint_document(file_path)
        elif suffix in ['.odt', '.ods', '.odp'] or 'opendocument' in mime_type:
            return self._extract_opendocument(file_path, suffix)
        else:
            raise ValueError(f"Unsupported Office document format: {suffix}")
    
    def _extract_word_document(self, file_path: Path) -> Dict[str, Any]:
        """Extract metadata from Word documents"""
        try:
            from docx import Document
            
            result = {}
            doc = Document(file_path)
            
            # Document properties
            props = doc.core_properties
            result['title'] = props.title
            result['author'] = props.author
            result['subject'] = props.subject
            result['creator'] = props.creator
            result['keywords'] = props.keywords.split(',') if props.keywords else []
            result['creation_date'] = self._parse_date(props.created)
            result['modification_date'] = self._parse_date(props.modified)
            result['language'] = props.language
            
            # Content analysis
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            result['text_content'] = '\n'.join(paragraphs)
            
            # Document structure
            result['page_count'] = self._estimate_page_count(result['text_content'])
            result['word_count'] = len(result['text_content'].split()) if result['text_content'] else 0
            result['character_count'] = len(result['text_content']) if result['text_content'] else 0
            
            # Count elements
            result['tables_count'] = len(doc.tables)
            result['images_count'] = self._count_images_in_word(doc)
            result['links_count'] = self._count_links_in_text(result['text_content'])
            
            # Additional metadata
            result['raw_metadata'] = {
                'category': props.category,
                'comments': props.comments,
                'content_status': props.content_status,
                'identifier': props.identifier,
                'revision': props.revision,
                'version': props.version
            }
            
            return result
            
        except ImportError:
            raise ImportError("python-docx library is required for Word document extraction")
        except Exception as e:
            self.logger.error(f"Word document extraction error: {e}")
            raise
    
    def _extract_excel_document(self, file_path: Path) -> Dict[str, Any]:
        """Extract metadata from Excel documents"""
        try:
            import openpyxl
            
            result = {}
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            
            # Document properties
            props = wb.properties
            result['title'] = props.title
            result['creator'] = props.creator
            result['subject'] = props.subject
            result['keywords'] = props.keywords.split(',') if props.keywords else []
            result['creation_date'] = self._parse_date(props.created)
            result['modification_date'] = self._parse_date(props.modified)
            
            # Workbook structure
            result['page_count'] = len(wb.worksheets)  # Sheets as "pages"
            
            # Extract text content from cells
            text_content = []
            total_cells = 0
            
            for worksheet in wb.worksheets:
                for row in worksheet.iter_rows(values_only=True):
                    for cell_value in row:
                        if cell_value is not None:
                            text_content.append(str(cell_value))
                            total_cells += 1
                            
                            # Limit extraction for performance
                            if total_cells > 10000:
                                break
                    if total_cells > 10000:
                        break
                if total_cells > 10000:
                    break
            
            result['text_content'] = ' '.join(text_content)
            result['word_count'] = len(result['text_content'].split()) if result['text_content'] else 0
            result['character_count'] = len(result['text_content']) if result['text_content'] else 0
            
            # Excel-specific metadata
            result['tables_count'] = len(wb.worksheets)  # Each sheet is essentially a table
            result['raw_metadata'] = {
                'sheet_count': len(wb.worksheets),
                'sheet_names': [ws.title for ws in wb.worksheets]
            }
            
            wb.close()
            return result
            
        except ImportError:
            raise ImportError("openpyxl library is required for Excel document extraction")
        except Exception as e:
            self.logger.error(f"Excel document extraction error: {e}")
            raise
    
    def _extract_powerpoint_document(self, file_path: Path) -> Dict[str, Any]:
        """Extract metadata from PowerPoint documents"""
        try:
            from pptx import Presentation
            
            result = {}
            prs = Presentation(file_path)
            
            # Document properties
            props = prs.core_properties
            result['title'] = props.title
            result['author'] = props.author
            result['subject'] = props.subject
            result['creator'] = props.creator
            result['keywords'] = props.keywords.split(',') if props.keywords else []
            result['creation_date'] = self._parse_date(props.created)
            result['modification_date'] = self._parse_date(props.modified)
            
            # Presentation structure
            result['page_count'] = len(prs.slides)
            
            # Extract text content from slides
            text_content = []
            image_count = 0
            
            for slide in prs.slides:
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text_content.append(shape.text)
                    # Count images
                    if shape.shape_type == 13:  # Picture type
                        image_count += 1
            
            result['text_content'] = '\n'.join(text_content)
            result['word_count'] = len(result['text_content'].split()) if result['text_content'] else 0
            result['character_count'] = len(result['text_content']) if result['text_content'] else 0
            result['images_count'] = image_count
            result['links_count'] = self._count_links_in_text(result['text_content'])
            
            # PowerPoint-specific metadata
            result['raw_metadata'] = {
                'slide_count': len(prs.slides)
            }
            
            return result
            
        except ImportError:
            raise ImportError("python-pptx library is required for PowerPoint document extraction")
        except Exception as e:
            self.logger.error(f"PowerPoint document extraction error: {e}")
            raise
    
    def _extract_opendocument(self, file_path: Path, suffix: str) -> Dict[str, Any]:
        """Extract metadata from OpenDocument files"""
        try:
            from odf import opendocument, text, table
            from odf.namespaces import OFFICENS
            
            result = {}
            
            if suffix == '.odt':
                doc = opendocument.load(file_path)
            elif suffix == '.ods':
                doc = opendocument.load(file_path)
            elif suffix == '.odp':
                doc = opendocument.load(file_path)
            else:
                raise ValueError(f"Unsupported OpenDocument format: {suffix}")
            
            # Extract metadata
            meta = doc.meta
            result['title'] = str(meta.getElementsByType(OFFICENS + 'title')[0]) if meta.getElementsByType(OFFICENS + 'title') else None
            result['creator'] = str(meta.getElementsByType(OFFICENS + 'creator')[0]) if meta.getElementsByType(OFFICENS + 'creator') else None
            result['subject'] = str(meta.getElementsByType(OFFICENS + 'subject')[0]) if meta.getElementsByType(OFFICENS + 'subject') else None
            
            # Extract text content
            text_elements = doc.getElementsByType(text.P)
            text_content = [str(element) for element in text_elements]
            result['text_content'] = '\n'.join(text_content)
            
            result['word_count'] = len(result['text_content'].split()) if result['text_content'] else 0
            result['character_count'] = len(result['text_content']) if result['text_content'] else 0
            
            # Count tables for ODS
            if suffix == '.ods':
                tables = doc.getElementsByType(table.Table)
                result['tables_count'] = len(tables)
            
            return result
            
        except ImportError:
            self.logger.warning("odfpy library not available for OpenDocument extraction")
            return {'text_content': '', 'extraction_method': 'fallback'}
        except Exception as e:
            self.logger.error(f"OpenDocument extraction error: {e}")
            raise
    
    def _count_images_in_word(self, doc) -> int:
        """Count images in Word document"""
        try:
            # This is a simplified count - full implementation would need more detailed analysis
            image_count = 0
            for paragraph in doc.paragraphs:
                for run in paragraph.runs:
                    if 'graphic' in run._r.xml:
                        image_count += 1
            return image_count
        except:
            return 0
    
    def _count_links_in_text(self, text: str) -> int:
        """Count URLs and email addresses in text"""
        if not text:
            return 0
        
        import re
        url_pattern = r'https?://[^\s]+'
        email_pattern = r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'
        
        urls = re.findall(url_pattern, text, re.IGNORECASE)
        emails = re.findall(email_pattern, text)
        
        return len(urls) + len(emails)
    
    def _estimate_page_count(self, text: str) -> int:
        """Estimate page count based on text length"""
        if not text:
            return 0
        
        # Rough estimation: 250 words per page
        words = len(text.split())
        estimated_pages = max(1, words // 250)
        return estimated_pages
    
    def _assess_extraction_quality(self, text: str, word_count: Optional[int]) -> float:
        """Assess quality of text extraction"""
        if not text:
            return 0.0
        
        # High confidence for Office documents as they have structured text
        confidence = 0.9
        
        # Adjust based on content length
        if word_count and word_count < 10:
            confidence *= 0.5
        elif word_count and word_count < 50:
            confidence *= 0.7
        
        return confidence
    
    def _suggest_categories(self, text: str, mime_type: str, office_data: Dict) -> List[str]:
        """Suggest categories based on document type and content"""
        categories = []
        
        # Document type based categories
        if 'wordprocessingml' in mime_type or '.doc' in str(mime_type):
            categories.append('document')
        elif 'spreadsheetml' in mime_type or '.xls' in str(mime_type):
            categories.append('spreadsheet')
        elif 'presentationml' in mime_type or '.ppt' in str(mime_type):
            categories.append('presentation')
        elif 'opendocument' in mime_type:
            categories.append('office-document')
        
        # Content-based categories
        if text:
            content_lower = text.lower()
            
            if any(word in content_lower for word in ['report', 'analysis', 'summary', 'findings']):
                categories.append('report')
            
            if any(word in content_lower for word in ['proposal', 'plan', 'strategy', 'recommendation']):
                categories.append('planning')
            
            if any(word in content_lower for word in ['budget', 'financial', 'cost', 'revenue', 'profit']):
                categories.append('financial')
            
            if any(word in content_lower for word in ['meeting', 'agenda', 'minutes', 'action items']):
                categories.append('meeting')
            
            if any(word in content_lower for word in ['manual', 'guide', 'instructions', 'procedure']):
                categories.append('documentation')
        
        return categories[:4]
    
    def _suggest_tags(self, text: str, title: Optional[str], mime_type: str) -> List[str]:
        """Suggest tags based on content, title, and document type"""
        tags = []
        
        # Document type tags
        if 'wordprocessingml' in mime_type:
            tags.extend(['word', 'document'])
        elif 'spreadsheetml' in mime_type:
            tags.extend(['excel', 'spreadsheet', 'data'])
        elif 'presentationml' in mime_type:
            tags.extend(['powerpoint', 'presentation', 'slides'])
        elif 'opendocument' in mime_type:
            tags.extend(['opendocument', 'office'])
        
        # Title-based tags
        if title:
            title_words = [word.strip('.,!?').lower() for word in title.split() if len(word) > 3]
            tags.extend(title_words[:5])
        
        # Content-based tags
        if text:
            content_tags = self._extract_keywords(text, max_keywords=8)
            tags.extend(content_tags)
        
        return list(set(tags))[:15]
    
    def get_supported_mime_types(self) -> List[str]:
        """Get list of supported MIME types"""
        return self.supported_types.copy()